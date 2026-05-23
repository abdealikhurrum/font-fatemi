#!/usr/bin/env python3
r"""
double_press_convert.py
Convert LSD double-press Arabic text to proper Unicode.

When typing on an LSD keyboard, secondary characters are entered by pressing
the same key twice quickly (سس → ے, جج → چھے, etc.).  Old documents may
contain these literal repeated-character sequences instead of the correct
Unicode characters.  This tool replaces them in-place.

Modes
-----
  python3 double_press_convert.py input.docx [output.docx]   # Word doc
  python3 double_press_convert.py input.pptx [output.pptx]   # PowerPoint
  python3 double_press_convert.py input.txt [output.txt]     # plain text file
  echo "سسضض" | python3 double_press_convert.py              # stdin → stdout
  cat file.txt | python3 double_press_convert.py

Options
-------
  -t, --tags TAGS         comma-separated tags to include
  -x, --exclude-tags TAGS comma-separated tags to exclude
  -r, --reverse           reverse conversion (Unicode → LSD sequences)

If output path is omitted the input file is overwritten.

"""

import sys
import re
import ast
import argparse
from typing import Sequence, Iterable, List, Tuple, Dict
from pathlib import Path

try:
    import docx
    from docx.oxml.ns import qn
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False


# ---------------------------------------------------------------------------
# Substitution tables
# ---------------------------------------------------------------------------

# LSD double-press: each entry is (repeated_char, secondary).
# Process longest-output entries first so جج→چھے doesn't become چ+extra.
_LSD_DOUBLES = [

    ("ض", "ٹ"), #double-daad
    ("ث", "پ"), #double-thaa
    ("ح", "چ"), #double-haa
    ("س", "ے"), #double-seen
    ("ك", "گ"), #double-kaaf
    ("ط", "ں"), #double-taa
    ("ظ", "ہ"), #double-zaa
]

_LSD_SINGLES = [
    ('\u0631'+'\u064C', '\u0691'),  #raa-toRraa ر + ً → ڑ
    ('\u062F'+'\u064C', '\u06BE'),  #dal-toDdaa د + ً → ڈ
    (' \u061B', '\u0686'+'\u06BE'+'\u06D2') #semicolon-toChe space + arabic semicolon → چھے
]

_ARABIC_WORD_CHARS = "\u0621-\u06FF"
_REVERSE_ISOLATED_CHEH_HEH_YEH = re.compile(rf"(?<![{_ARABIC_WORD_CHARS}])\u0686\u06BE\u06D2(?![{_ARABIC_WORD_CHARS}])")
_REVERSE_YEH_WORD_END = re.compile(rf"\u06D2(?=(?:[^{_ARABIC_WORD_CHARS}]|$))")


def _parse_annotated_substitutions() -> Tuple[List[Tuple[str, str, List[str]]], List[Tuple[str, str, List[str]]]]:
    """Parse this source file to extract tags from end-of-line comments.

    Returns two lists: doubles and singles. Each item is (src, dst, [tags]).
    """
    src_path = Path(__file__)
    src_text = src_path.read_text(encoding="utf8")

    def _extract(block_name: str) -> List[Tuple[str, str, List[str]]]:
        m = re.search(rf"{block_name}\s*=\s*\[", src_text)
        if not m:
            return []
        start = m.end()
        # find the closing bracket for the list
        end = src_text.find("]", start)
        block = src_text[start:end]
        items: List[Tuple[str, str, List[str]]] = []
        for line in block.splitlines():
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            code_part, _, comment = line.partition("#")
            code_part = code_part.rstrip(", ")
            try:
                pair = ast.literal_eval(code_part)
            except Exception:
                continue
            tags = []
            if comment:
                # tags are words in the comment (split on spaces/commas)
                tags = [t.strip().lstrip('#') for t in re.split(r"[\s,]+", comment.strip()) if t.strip()]
            items.append((pair[0], pair[1], tags))
        return items

    doubles = _extract("_LSD_DOUBLES")
    singles = _extract("_LSD_SINGLES")
    return doubles, singles


def _apply_tag_filter(include: Iterable[str] = None, exclude: Iterable[str] = None):
    """Set active substitutions based on include/exclude tag lists.

    If `include` is provided, only substitutions with any matching tag are kept.
    If `exclude` is provided, substitutions with any matching tag are removed.
    """
    global _ACTIVE_DOUBLES, _ACTIVE_SINGLES, _lsd_replace, _lsd_singles_replace

    include_set = set(include) if include else None
    exclude_set = set(exclude) if exclude else set()

    parsed_doubles, parsed_singles = _parse_annotated_substitutions()

    def _filter(items):
        out = []
        for src, dst, tags in items:
            tagset = set(tags)
            if include_set is not None:
                if not (tagset & include_set):
                    continue
            if tagset & exclude_set:
                continue
            out.append((src, dst))
        return out

    _ACTIVE_DOUBLES = _filter(parsed_doubles) or _LSD_DOUBLES[:]
    _ACTIVE_SINGLES = _filter(parsed_singles) or _LSD_SINGLES[:]

    _lsd_replace = _make_replacer(_ACTIVE_DOUBLES)
    _lsd_singles_replace = _make_singles_replacer(_ACTIVE_SINGLES)
    # rebuild reverse replacers when active lists change
    global _rev_double_replace, _rev_single_replace, _REV_DOUBLE_LOOKUP, _REV_SINGLE_LOOKUP
    _rev_double_replace, _rev_single_replace, _REV_DOUBLE_LOOKUP, _REV_SINGLE_LOOKUP = _make_reverse_replacers(_ACTIVE_DOUBLES, _ACTIVE_SINGLES)


def _make_reverse_replacers(doubles, singles):
    """Create reverse replacers mapping outputs back to original sequences."""
    # For doubles: dst -> src+src, excluding yeh because it's context-sensitive.
    rev_doubles = [(dst, c + c) for c, dst in doubles if dst != '\u06D2']
    # For singles: dst -> src, excluding cheh_heh_yeh because it's context-sensitive.
    rev_singles = [(dst, src) for src, dst in singles if dst != '\u0686\u06BE\u06D2']

    def _build_pattern(items):
        # longest-first to avoid partial matches
        alts = sorted([re.escape(a) for a, _ in items], key=len, reverse=True)
        if not alts:
            return re.compile("$")
        pattern = "|".join(alts)
        # only reverse these doubles at a word boundary/end to avoid mis-reverting normal text
        return re.compile(rf"(?:{pattern})(?=(?:[^{_ARABIC_WORD_CHARS}]|$))")

    def _make(items):
        lookup = {a: b for a, b in items}
        pattern = _build_pattern(items)

        def replace(text: str) -> str:
            return pattern.sub(lambda m: lookup[m.group(0)], text)

        return replace, lookup

    rev_d_replacer, rev_d_lookup = _make(rev_doubles)
    rev_s_replacer, rev_s_lookup = _make(rev_singles)
    return rev_d_replacer, rev_s_replacer, rev_d_lookup, rev_s_lookup

# Initialize active lists to defaults (all substitutions enabled)
_ACTIVE_DOUBLES = _LSD_DOUBLES[:]
_ACTIVE_SINGLES = _LSD_SINGLES[:]

# initial replacers (created after helper functions below)
def _build_pattern(doubles):
    """Compile a single regex that matches all double-press pairs."""
    alts = [re.escape(c + c) for c, _ in doubles]
    return re.compile("|".join(alts))


def _make_replacer(doubles):
    lookup = {c + c: s for c, s in doubles}
    pattern = _build_pattern(doubles)

    def replace(text: str) -> str:
        return pattern.sub(lambda m: lookup[m.group()], text)

    return replace


_lsd_replace = _make_replacer(_ACTIVE_DOUBLES)


def _build_singles_pattern(singles):
    alts = [re.escape(s) for s, _ in singles]
    return re.compile("|".join(alts))


def _make_singles_replacer(singles):
    lookup = {s: r for s, r in singles}
    pattern = _build_singles_pattern(singles)

    def replace(text: str) -> str:
        return pattern.sub(lambda m: lookup[m.group()], text)

    return replace


_lsd_singles_replace = _make_singles_replacer(_ACTIVE_SINGLES)

# initialize reverse replacers based on active lists
_rev_double_replace, _rev_single_replace, _REV_DOUBLE_LOOKUP, _REV_SINGLE_LOOKUP = _make_reverse_replacers(_ACTIVE_DOUBLES, _ACTIVE_SINGLES)


def convert_text(text: str, crulp: bool = False) -> str:
    """Apply double-press substitutions to a plain string."""
    # double-press replacements
    text = _lsd_replace(text)
    # single-sequence replacements (e.g. char+diacritic → single glyph)
    text = _lsd_singles_replace(text)
    return text


def convert_text_reverse(text: str) -> str:
    """Reverse substitutions applied by convert_text."""
    # Apply context-sensitive reverse rules first.
    text = _REVERSE_ISOLATED_CHEH_HEH_YEH.sub(' \u061B', text)
    text = _REVERSE_YEH_WORD_END.sub('سس', text)
    # then apply generic reverse mappings for remaining substitutions.
    text = _rev_single_replace(text)
    text = _rev_double_replace(text)
    return text


# ---------------------------------------------------------------------------
# Word document processing
# ---------------------------------------------------------------------------

def _all_paragraphs(doc):
    """Yield every paragraph in the document (body + tables + headers/footers)."""
    yield from doc.paragraphs
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs
    for section in doc.sections:
        for hf in (section.header, section.footer,
                   section.even_page_header, section.even_page_footer,
                   section.first_page_header, section.first_page_footer):
            if hf is not None:
                yield from hf.paragraphs


def _fix_paragraph(para):
    """
    Apply substitutions within each run, then handle cross-run boundaries.

    Cross-run pass: if run[i] ends with X and run[i+1] starts with X and
    X+X is a double-press pair, merge those two characters and substitute.
    """
    runs = para.runs
    if not runs:
        return

    # Within-run substitution
    for run in runs:
        if run.text:
            run.text = convert_text(run.text)

    # Cross-run boundary pass (one pass is sufficient since substitution
    # consumes both copies of the character)
    doubles = _ACTIVE_DOUBLES
    lookup = {c + c: s for c, s in doubles}

    for i in range(len(runs) - 1):
        a, b = runs[i], runs[i + 1]
        if not a.text or not b.text:
            continue
        pair = a.text[-1] + b.text[0]
        if pair in lookup:
            replacement = lookup[pair]
            a.text = a.text[:-1] + replacement
            b.text = b.text[1:]

    # Cross-run singles pass: handle sequences split across runs
    singles = _ACTIVE_SINGLES
    sing_lookup = {s: r for s, r in singles}
    for i in range(len(runs) - 1):
        a, b = runs[i], runs[i + 1]
        if not a.text or not b.text:
            continue
        pair = a.text[-1] + b.text[0]
        if pair in sing_lookup:
            replacement = sing_lookup[pair]
            a.text = a.text[:-1] + replacement
            b.text = b.text[1:]



def convert_docx(src_path: str, dst_path: str):
    doc = docx.Document(src_path)
    changed = 0
    for para in _all_paragraphs(doc):
        before = "".join(r.text for r in para.runs)
        _fix_paragraph(para)
        after = "".join(r.text for r in para.runs)
        if before != after:
            changed += 1
    doc.save(dst_path)
    return changed


def _fix_paragraph_reverse(para):
    runs = para.runs
    if not runs:
        return

    # Within-run reverse substitution
    for run in runs:
        if run.text:
            run.text = convert_text_reverse(run.text)

    # Cross-run boundary reverse pass
    for i in range(len(runs) - 1):
        a, b = runs[i], runs[i + 1]
        if not a.text or not b.text:
            continue
        for key, repl in {**_REV_DOUBLE_LOOKUP, **_REV_SINGLE_LOOKUP}.items():
            L = len(key)
            for k in range(1, L):
                if a.text.endswith(key[:k]) and b.text.startswith(key[k:]):
                    a.text = a.text[:-k] + repl
                    b.text = b.text[L-k:]
                    break


def convert_docx_reverse(src_path: str, dst_path: str):
    doc = docx.Document(src_path)
    changed = 0
    for para in _all_paragraphs(doc):
        before = "".join(r.text for r in para.runs)
        _fix_paragraph_reverse(para)
        after = "".join(r.text for r in para.runs)
        if before != after:
            changed += 1
    doc.save(dst_path)
    return changed


def convert_txt(src_path: str, dst_path: str):
    """Convert a plain text file in-place or to a new file."""
    text = Path(src_path).read_text(encoding="utf8")
    out = convert_text(text)
    changed = 1 if out != text else 0
    Path(dst_path).write_text(out, encoding="utf8")
    return changed


def convert_txt_reverse(src_path: str, dst_path: str):
    text = Path(src_path).read_text(encoding="utf8")
    out = convert_text_reverse(text)
    changed = 1 if out != text else 0
    Path(dst_path).write_text(out, encoding="utf8")
    return changed


def convert_pptx(src_path: str, dst_path: str):
    """Convert text in a .pptx Presentation (all text-containing shapes)."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(src_path)
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            try:
                text_frame = shape.text_frame
            except Exception:
                # some shape types (group shapes, legacy placeholders) may
                # raise when accessing text_frame; skip them safely.
                continue
            if text_frame is None:
                continue
            for para in text_frame.paragraphs:
                before = "".join(r.text for r in para.runs)
                # Within-run replacements
                for run in para.runs:
                    if run.text:
                        run.text = convert_text(run.text)

                # Cross-run boundary passes (doubles and singles)
                doubles = _ACTIVE_DOUBLES
                lookup = {c + c: s for c, s in doubles}
                for i in range(len(para.runs) - 1):
                    a, b = para.runs[i], para.runs[i + 1]
                    if not a.text or not b.text:
                        continue
                    pair = a.text[-1] + b.text[0]
                    if pair in lookup:
                        replacement = lookup[pair]
                        a.text = a.text[:-1] + replacement
                        b.text = b.text[1:]

                singles = _ACTIVE_SINGLES
                sing_lookup = {s: r for s, r in singles}
                for i in range(len(para.runs) - 1):
                    a, b = para.runs[i], para.runs[i + 1]
                    if not a.text or not b.text:
                        continue
                    pair = a.text[-1] + b.text[0]
                    if pair in sing_lookup:
                        replacement = sing_lookup[pair]
                        a.text = a.text[:-1] + replacement
                        b.text = b.text[1:]

                after = "".join(r.text for r in para.runs)
                if before != after:
                    changed += 1
    prs.save(dst_path)
    return changed


def convert_pptx_reverse(src_path: str, dst_path: str):
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(src_path)
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            try:
                text_frame = shape.text_frame
            except Exception:
                continue
            if text_frame is None:
                continue
            for para in text_frame.paragraphs:
                before = "".join(r.text for r in para.runs)
                # Within-run reverse replacements
                for run in para.runs:
                    if run.text:
                        run.text = convert_text_reverse(run.text)

                # Cross-run boundary reverse passes
                # handle reverse-doubles (dst -> src+src)
                for i in range(len(para.runs) - 1):
                    a, b = para.runs[i], para.runs[i + 1]
                    if not a.text or not b.text:
                        continue
                    # try all reverse-double keys split across runs
                    for key, repl in _REV_DOUBLE_LOOKUP.items():
                        L = len(key)
                        for k in range(1, L):
                            if a.text.endswith(key[:k]) and b.text.startswith(key[k:]):
                                a.text = a.text[:-k] + repl
                                b.text = b.text[L-k:]
                                break

                # handle reverse-singles
                for i in range(len(para.runs) - 1):
                    a, b = para.runs[i], para.runs[i + 1]
                    if not a.text or not b.text:
                        continue
                    for key, repl in _REV_SINGLE_LOOKUP.items():
                        L = len(key)
                        for k in range(1, L):
                            if a.text.endswith(key[:k]) and b.text.startswith(key[k:]):
                                a.text = a.text[:-k] + repl
                                b.text = b.text[L-k:]
                                break

                after = "".join(r.text for r in para.runs)
                if before != after:
                    changed += 1
    prs.save(dst_path)
    return changed


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv: Sequence[str]):
    parser = argparse.ArgumentParser(prog=argv[0],
                                     description='Convert LSD double-press Arabic text to proper Unicode.')
    parser.add_argument('-t', '--tags', help='Comma-separated list of tags to include')
    parser.add_argument('-x', '--exclude-tags', help='Comma-separated list of tags to exclude')
    parser.add_argument('-r', '--reverse', action='store_true', help='Run reverse conversion (Unicode -> LSD sequences)')
    parser.add_argument('src', nargs='?', help='Source .docx file or omitted for stdin')
    parser.add_argument('dst', nargs='?', help='Destination .docx file (optional)')

    ns = parser.parse_args(argv[1:])

    include = ns.tags.split(',') if ns.tags else None
    exclude = ns.exclude_tags.split(',') if ns.exclude_tags else None
    reverse = bool(ns.reverse)
    # apply tag filtering (no-op if include/exclude are None)
    _apply_tag_filter(include, exclude)

    if not ns.src and not sys.stdin.isatty():
        # Plain-text filter mode
        text = sys.stdin.read()
        sys.stdout.write(convert_text(text))
        return

    if not ns.src:
        print(__doc__)
        sys.exit(0)

    src = ns.src
    dst = ns.dst if ns.dst else src

    # Dispatch by extension
    lower = src.lower()
    if lower.endswith('.docx'):
        if not DOCX_AVAILABLE:
            print("Error: python-docx not installed.  Run: pip install python-docx", file=sys.stderr)
            sys.exit(1)
        if reverse:
            n = convert_docx_reverse(src, dst)
        else:
            n = convert_docx(src, dst)
    elif lower.endswith('.pptx'):
        if not PPTX_AVAILABLE:
            print("Error: python-pptx not installed.  Run: pip install python-pptx", file=sys.stderr)
            sys.exit(1)
        if reverse:
            n = convert_pptx_reverse(src, dst)
        else:
            n = convert_pptx(src, dst)
    elif lower.endswith('.txt'):
        if reverse:
            n = convert_txt_reverse(src, dst)
        else:
            n = convert_txt(src, dst)
    else:
        print(f"Error: unsupported file type: {src}", file=sys.stderr)
        sys.exit(1)

    layout = "LSD"
    action = "overwritten" if dst == src else f"saved to {dst!r}"
    print(f"{layout}: {n} paragraph(s) changed — {action}")

    layout = "LSD"
    action = "overwritten" if dst == src else f"saved to {dst!r}"
    print(f"{layout}: {n} paragraph(s) changed — {action}")


if __name__ == "__main__":
    main(sys.argv)
